"""
Build a CTO-facing PPTX deck from a data.json artefact.

Source-agnostic. Reads `source`, `window_label`, `summary`, `dims` from
the JSON and generates prose dynamically based on the observed numbers.
Slides that depend on a missing dimension (e.g. no principal field on
the source) are automatically skipped.

Usage:
    python3 build_pptx.py --data reports/Prompt_Security_7d/data.json

Output lands next to the input JSON
(reports/<slug>_<window>/<slug>_CTO_Deck_<window>.pptx).
"""
from __future__ import annotations

import argparse
import json
import re
from collections import defaultdict
from datetime import datetime, timezone
from pathlib import Path
from typing import Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# palette (matches DOCX + charts)
NAVY = RGBColor(0x0B, 0x1F, 0x3A)
CORAL = RGBColor(0xE8, 0x43, 0x4F)
AMBER = RGBColor(0xF2, 0xA9, 0x50)
TEAL = RGBColor(0x4F, 0xB8, 0x9E)
SLATE = RGBColor(0x88, 0x99, 0xA6)
CREAM = RGBColor(0xF5, 0xF1, 0xE8)
INK = RGBColor(0x1A, 0x1F, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x55, 0x5F, 0x70)
LIGHT = RGBColor(0xCA, 0xD3, 0xE0)
GRAY = RGBColor(0x9C, 0xA9, 0xBB)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _slugify(name: str) -> str:
    s = re.sub(r"[^0-9A-Za-z]+", "_", name.strip())
    return re.sub(r"_+", "_", s).strip("_") or "source"


# ------------------------------------------------------- pptx helpers

def set_slide_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    spTree = slide.shapes._spTree
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def add_text(slide, left, top, width, height, text,
             font="Calibri", size=14, bold=False, italic=False,
             color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.15):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def add_accent_bar(slide, left, top, width, height, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def add_footer(slide, page_num, total, source, window, light_bg=True):
    fg = MUTED if light_bg else GRAY
    add_text(slide,
             Inches(0.5), Inches(7.15),
             Inches(8), Inches(0.25),
             f"SentinelOne + {source}  |  {window} CTO review",
             size=9, color=fg)
    add_text(slide,
             Inches(11.0), Inches(7.15),
             Inches(1.8), Inches(0.25),
             f"{page_num} / {total}",
             size=9, color=fg, align=PP_ALIGN.RIGHT)


def slide_title_bar(slide, title, eyebrow=None):
    add_accent_bar(slide, 0, 0, SLIDE_W, Inches(0.12), NAVY)
    if eyebrow:
        add_text(slide,
                 Inches(0.5), Inches(0.28),
                 Inches(12), Inches(0.28),
                 eyebrow.upper(),
                 size=10, bold=True, color=CORAL)
    add_text(slide,
             Inches(0.5), Inches(0.55),
             Inches(12.3), Inches(0.6),
             title,
             size=28, bold=True, color=NAVY)


# ------------------------------------------------------- commentary

def _mix_rows(d):
    return (d["queries"].get("per_user_mix_top10", {})
            .get("rows", []) or [])


def _principal_key(mix_rows) -> Optional[str]:
    for r in mix_rows:
        for k in ("user", "src.hostname", "src.ip.address"):
            if k in r:
                return k
    return None


def _principal_label(p_key: Optional[str]) -> str:
    if not p_key:
        return "principals"
    return {"user": "users", "src.hostname": "hosts",
            "src.ip.address": "source IPs"}.get(p_key, p_key)


def _derive_by_principal(mix_rows, p_key):
    acc = defaultdict(int)
    for r in mix_rows:
        who = r.get(p_key)
        if who in (None, "", "null"):
            continue
        acc[str(who)] += int(r.get("n") or 0)
    return sorted(acc.items(), key=lambda kv: -kv[1])


def _fmt_pct(p: float) -> str:
    if p >= 10:
        return f"{p:.1f}%"
    if p >= 1:
        return f"{p:.2f}%"
    return f"{p:.3f}%"


def _concentration_short(top_share: float, who: str) -> str:
    if top_share >= 90:
        return (f"{who} drove {top_share:.1f}% of all events. "
                "Resolve this account before treating any other number "
                "as a baseline.")
    if top_share >= 50:
        return (f"{who} alone drove {top_share:.1f}% of volume. Tag "
                "this principal before deriving steady-state metrics.")
    if top_share >= 20:
        return (f"{who} leads the window at {top_share:.1f}% of events. "
                "Not dominant, but worth attention.")
    return (f"No single principal dominates. {who} leads at "
            f"{top_share:.1f}% of the window.")


def _intervention_short(pct: float) -> str:
    if pct >= 40:
        return (f"Intervention rate is {pct:.1f}%. The policy is doing "
                "heavy work. Confirm the traffic mix justifies this "
                "posture before moving anything to enforce.")
    if pct >= 10:
        return (f"Intervention rate is {pct:.1f}%. Policy is active on "
                "a meaningful slice without blocking the business.")
    return (f"Intervention rate is {pct:.1f}%. The source is in "
            "observe posture; most events pass through.")


def _bypass_short(pct: float) -> str:
    if pct < 0.5:
        return (f"Bypass is near-zero ({pct:.2f}%). Users are accepting "
                "policy decisions. The cleanest possible trust signal.")
    if pct < 5:
        return (f"Bypass sits at {pct:.2f}%. Low-volume but "
                "high-signal: these are principals disagreeing with "
                "the policy.")
    return (f"Bypass is elevated at {pct:.2f}%. A meaningful share of "
            "traffic is opting out; review which policies are involved.")


def _recommendations(d) -> List[Tuple[str, RGBColor, str, str]]:
    summary = d["summary"]
    dims = d.get("dims", {})
    p_key = summary.get("top_principal_key")
    tu = summary.get("top_user")
    source = d["source"]
    recs: List[Tuple[str, RGBColor, str, str]] = []

    if tu and summary.get("top_share", 0.0) >= 40:
        who = tu.get(p_key) if p_key else "top principal"
        recs.append((
            "24 HOURS", CORAL,
            f"Resolve {who}",
            f"A single {_principal_label(p_key).rstrip('s')} drove "
            f"{summary['top_share']:.1f}% of the window. Pivot into "
            "EDR / identity from that session to establish harness vs. "
            "shared-account vs. abuse. Close the single largest risk "
            "in this report."))

    if (dims.get("action")
            and summary.get("block_pct", 0.0) >= 10
            and summary.get("bypass_pct", 0.0) < 1.0):
        recs.append((
            "THIS WEEK", AMBER,
            "Promote log-only policies to enforce",
            f"Block {_fmt_pct(summary['block_pct'])} and bypass "
            f"{_fmt_pct(summary['bypass_pct'])}. That combination is "
            "as clean a trust signal as we get. Walk each log-only "
            "policy and flip the ones with zero false positives."))

    if dims.get("action") and summary.get("bypass", 0) > 0:
        recs.append((
            "THIS MONTH", TEAL,
            "Correlate bypass with endpoint posture",
            f"STAR rule: {source} bypass THEN a SentinelOne alert of "
            "severity medium+ on the same principal within 60 minutes. "
            "One rule, two data sources, one unified alert."))

    recs.append((
        "THIS MONTH", NAVY,
        "Enrich ingestion where it is shallow",
        f"Today, {source} lands with a subset of fields populated. "
        "Add prompt / payload / content where possible; that unlocks "
        "content-based detections and Purple AI hunts across the same "
        "data lake."))

    recs.append((
        "QUARTERLY", SLATE,
        "Baseline normal",
        "Re-run this report on a fixed cadence. Charts regenerate from "
        "a single JSON artefact, so 30- or 90-day views are a "
        "parameter flip. Anomalies feed the next round of tuning."))

    return recs[:4]


# ------------------------------------------------------- slides

def slide_cover(prs, d):
    source = d["source"]
    window = d["window_label"]
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(s, NAVY)

    add_accent_bar(s, 0, Inches(2.4), Inches(0.4), Inches(3.2), CORAL)

    add_text(s, Inches(0.9), Inches(2.35),
             Inches(11), Inches(0.4),
             f"CTO REPORT  |  {source.upper()}",
             size=14, bold=True, color=CORAL)

    add_text(s, Inches(0.9), Inches(2.8),
             Inches(11.5), Inches(1.5),
             f"{window} of activity\nunder policy",
             size=52, bold=True, color=WHITE, line_spacing=1.0)

    add_text(s, Inches(0.9), Inches(5.1),
             Inches(11), Inches(0.6),
             f"What {source} observed, what we intervened on,\n"
             "and where SentinelOne takes the insight next.",
             size=18, italic=True, color=LIGHT, line_spacing=1.25)

    add_text(s, Inches(0.9), Inches(6.6),
             Inches(7), Inches(0.3),
             f"Window   {d['window_start'][:10]}  to  "
             f"{d['window_end'][:10]}  UTC",
             size=11, color=LIGHT)
    add_text(s, Inches(0.9), Inches(6.9),
             Inches(8), Inches(0.3),
             f"Data source   dataSource.name = '{source}'",
             size=11, color=LIGHT, font="Consolas")
    add_text(s, Inches(9.5), Inches(6.75),
             Inches(3.5), Inches(0.3),
             "SentinelOne Pre-Sales Engineering",
             size=10, color=GRAY, align=PP_ALIGN.RIGHT)
    add_text(s, Inches(9.5), Inches(7.0),
             Inches(3.5), Inches(0.3),
             datetime.now(timezone.utc).strftime("Generated %Y-%m-%d"),
             size=10, color=GRAY, align=PP_ALIGN.RIGHT)


def stat_card(slide, left, top, width, height, value, label,
              accent=NAVY, unit=""):
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xDD, 0xE1, 0xE6)
    card.line.width = Pt(0.5)
    card.shadow.inherit = False
    add_accent_bar(slide, left, top, Inches(0.12), height, accent)
    val_txt = f"{value}{unit}" if unit else str(value)
    # shrink large-value font when the string is long so it fits on
    # one line inside a 2.95" card (numbers stay at 40pt)
    n = len(val_txt)
    if n <= 7:
        val_size = 40
    elif n <= 12:
        val_size = 28
    elif n <= 18:
        val_size = 20
    else:
        val_size = 16
    add_text(slide, left + Inches(0.3), top + Inches(0.2),
             width - Inches(0.45), Inches(1.1),
             val_txt, size=val_size, bold=True, color=NAVY,
             line_spacing=1.0)
    add_text(slide, left + Inches(0.3),
             top + height - Inches(0.5),
             width - Inches(0.45), Inches(0.4),
             label, size=11, color=MUTED)


def slide_execsummary(prs, d):
    summary = d["summary"]
    source = d["source"]
    window = d["window_label"]
    dims = d.get("dims", {})
    mix = _mix_rows(d)
    p_key = summary.get("top_principal_key") or _principal_key(mix)
    principals = _derive_by_principal(mix, p_key) if p_key else []
    principal_count = len(principals)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title_bar(s, "Executive summary",
                    eyebrow=f"{window} snapshot")

    # 4 stat cards
    card_w = Inches(2.95)
    card_h = Inches(1.6)
    top = Inches(1.4)
    gap = Inches(0.15)
    left0 = Inches(0.5)

    stat_card(s, left0 + 0 * (card_w + gap), top, card_w, card_h,
              f"{summary['total']:,}", f"Total events ({window})",
              accent=NAVY)
    if dims.get("action"):
        stat_card(s, left0 + 1 * (card_w + gap), top, card_w, card_h,
                  f"{summary['block']:,}", "Blocked events",
                  accent=CORAL)
        stat_card(s, left0 + 2 * (card_w + gap), top, card_w, card_h,
                  f"{summary['intervention_pct']:.1f}",
                  "Policy intervention rate",
                  accent=AMBER, unit="%")
    else:
        # no action field: show the dominant primary key + slice count
        # instead of stranded "n/a" cards
        by = summary.get("by_action") or {}
        rows = sorted(by.items(), key=lambda kv: -kv[1])
        dom_label = rows[0][0] if rows else "(none)"
        if len(dom_label) > 22:
            dom_label = dom_label[:19] + "..."
        stat_card(s, left0 + 1 * (card_w + gap), top, card_w, card_h,
                  dom_label,
                  f"Dominant {summary.get('prim_key', 'value')}",
                  accent=SLATE)
        n_slices = d.get("strategy", {}).get("n_slices", "?")
        stat_card(s, left0 + 2 * (card_w + gap), top, card_w, card_h,
                  f"{n_slices}", "Timeline slices",
                  accent=AMBER)
    if p_key and principal_count > 0:
        stat_card(s, left0 + 3 * (card_w + gap), top, card_w, card_h,
                  f"{principal_count}",
                  f"Distinct {_principal_label(p_key)}",
                  accent=TEAL)
    else:
        rank = summary.get("rank_24h")
        stat_card(s, left0 + 3 * (card_w + gap), top, card_w, card_h,
                  f"#{rank}" if rank else "n/a",
                  "Tenant rank (24h)",
                  accent=TEAL)

    # three takeaways
    add_text(s, Inches(0.5), Inches(3.3),
             Inches(12.3), Inches(0.4),
             "Three things for the CTO to act on",
             size=18, bold=True, color=NAVY)

    items: List[Tuple[str, str]] = []
    if dims.get("action"):
        items.append((
            f"Block {_fmt_pct(summary['block_pct'])}, bypass "
            f"{_fmt_pct(summary['bypass_pct'])}.",
            _intervention_short(summary.get("intervention_pct", 0.0))))
    if summary.get("top_user") and p_key:
        tu = summary["top_user"]
        items.append((
            f"One {_principal_label(p_key).rstrip('s')} drove "
            f"{summary['top_share']:.1f}% of volume.",
            _concentration_short(summary["top_share"], tu.get(p_key))))
    if dims.get("action") and summary.get("bypass_pct") is not None:
        items.append((
            "Bypass is the signal to watch.",
            _bypass_short(summary.get("bypass_pct", 0.0))))

    # fallbacks for sources without action/principal
    if not items:
        prim = summary.get("prim_key")
        by = summary.get("by_action") or {}
        rows = sorted(by.items(), key=lambda kv: -kv[1])
        if rows and prim:
            top_label, top_n = rows[0]
            share = (100.0 * top_n / summary["total"]
                     if summary["total"] else 0)
            items.append((
                f"Dominant {prim}: '{top_label}'.",
                f"{top_n:,} events ({share:.1f}% of the window). Use "
                "this as the entry point when building detections or "
                "hunts against the source."))
        if summary.get("rank_24h"):
            items.append((
                f"Tenant rank #{summary['rank_24h']} in the last 24h.",
                f"{source} sits in the top tier of data sources landing "
                "in this tenant. A useful anchor for whether volume is "
                "trending or flat."))
        items.append((
            "Same data lake as EDR + identity + network.",
            f"Any detection or hunt built on {source} can join to "
            "SentinelOne endpoint context in one PowerQuery. No "
            "cross-system plumbing."))

    while len(items) < 3:
        items.append((
            "Same lake as EDR / identity / network.",
            "Every chart in this deck joins natively with SentinelOne "
            "endpoint context in a single PowerQuery."))

    y = Inches(3.85)
    for title, body in items[:3]:
        add_accent_bar(s, Inches(0.5), y + Inches(0.08),
                       Inches(0.1), Inches(0.7), CORAL)
        add_text(s, Inches(0.75), y,
                 Inches(12.1), Inches(0.35),
                 title, size=14, bold=True, color=NAVY)
        add_text(s, Inches(0.75), y + Inches(0.38),
                 Inches(12.1), Inches(0.55),
                 body, size=11.5, color=INK, line_spacing=1.2)
        y += Inches(0.95)


def slide_chart(prs, title, eyebrow, chart_path, takeaways, accent=NAVY):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title_bar(s, title, eyebrow=eyebrow)
    if chart_path.exists():
        s.shapes.add_picture(str(chart_path),
                             Inches(0.4), Inches(1.5),
                             width=Inches(7.8))
    else:
        add_text(s, Inches(0.4), Inches(3.5),
                 Inches(7.8), Inches(0.4),
                 f"(chart missing: {chart_path.name})",
                 size=14, italic=True, color=CORAL)
    panel_left = Inches(8.4)
    panel_top = Inches(1.5)
    panel_w = Inches(4.5)
    panel_h = Inches(5.2)
    panel = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               panel_left, panel_top, panel_w, panel_h)
    panel.fill.solid()
    panel.fill.fore_color.rgb = CREAM
    panel.line.fill.background()
    panel.shadow.inherit = False
    add_accent_bar(s, panel_left, panel_top, panel_w, Inches(0.1), accent)

    add_text(s, panel_left + Inches(0.25), panel_top + Inches(0.25),
             panel_w - Inches(0.5), Inches(0.3),
             "What to notice", size=11, bold=True, color=accent)
    y = panel_top + Inches(0.7)
    for t in takeaways:
        add_text(s, panel_left + Inches(0.25), y,
                 Inches(0.2), Inches(0.3),
                 "\u25A0", size=10, bold=True, color=accent)
        add_text(s, panel_left + Inches(0.5), y,
                 panel_w - Inches(0.75), Inches(1.4),
                 t, size=11, color=INK, line_spacing=1.25)
        y += Inches(1.1)
    return s


def slide_action_mix(prs, d, charts_dir):
    if not d.get("dims", {}).get("action"):
        return False
    summary = d["summary"]
    by = {str(r.get("action")): int(r.get("n") or 0)
          for r in d["queries"]["by_action"]["rows"]}
    total = summary["total"]
    block_p = 100.0 * by.get("block", 0) / total if total else 0
    modify_p = 100.0 * by.get("modify", 0) / total if total else 0
    bypass_p = 100.0 * by.get("bypass", 0) / total if total else 0

    takeaways = []
    if by.get("block"):
        takeaways.append(
            f"BLOCK at {_fmt_pct(block_p)}. Events the policy refused "
            "outright; exfil, compliance, jailbreak attempts live here.")
    if by.get("modify"):
        takeaways.append(
            f"MODIFY at {_fmt_pct(modify_p)}. Payload rewritten in "
            "place (eg. PII redacted) without stopping the user.")
    if by.get("bypass") is not None:
        takeaways.append(
            f"BYPASS at {_fmt_pct(bypass_p)}. "
            + ("Users almost never override. Policy is trusted."
               if bypass_p < 0.5 else
               "Override rate worth monitoring for drift."))
    # pad if we only got one or two
    if len(takeaways) < 3 and by.get("log"):
        log_p = 100.0 * by.get("log", 0) / total if total else 0
        takeaways.append(
            f"LOG at {_fmt_pct(log_p)}. Passed through and recorded "
            "for audit; the denominator everything else divides against.")

    slide_chart(prs, "Policy action mix", eyebrow="section 1",
                chart_path=charts_dir / "01_action_mix.png",
                takeaways=takeaways[:3], accent=CORAL)
    return True


def slide_daily_timeline(prs, d, charts_dir):
    source = d["source"]
    rows = (d["queries"].get("daily_by_action", {})
            .get("rows", []) or [])
    totals = [int(r.get("matchCount") or 0) for r in rows]

    takeaways: List[str] = []
    if totals:
        mx = max(totals)
        mx_idx = totals.index(mx)
        mx_label = rows[mx_idx].get("date", "?")
        tot_sum = sum(totals) or 1
        share = 100 * mx / tot_sum
        if share >= 70:
            takeaways.append(
                f"Heavily skewed: {mx_label} alone carried "
                f"{share:.1f}% of the window ({mx:,} events). "
                "Tag the spike out before calculating baselines.")
        elif share >= 30:
            takeaways.append(
                f"Uneven but not pathological: busiest slice "
                f"({mx_label}) carried {share:.1f}% of the window.")
        else:
            takeaways.append(
                f"Evenly distributed across the window. Busiest slice "
                f"({mx_label}) only {share:.1f}% ahead of average.")
    else:
        takeaways.append(
            "Timeline unavailable: no slices returned events in the "
            "window.")

    # stability note adapts to whether we have an action field
    has_action = d.get("dims", {}).get("action")
    prim = d.get("summary", {}).get("prim_key") or "primary key"
    noun = "action" if has_action else prim
    action_tot = defaultdict(int)
    for r in rows:
        for k, v in (r.get("by_action", {}) or {}).items():
            action_tot[k] += int(v)
    if action_tot and has_action:
        primary = max(action_tot.items(), key=lambda kv: kv[1])[0]
        takeaways.append(
            f"Dominant {noun} across the window: '{primary}'. Watch "
            f"the ratio to the next {noun} as an early policy-drift "
            "signal.")
        takeaways.append(
            "Policy decisions are stable under load across the window. "
            "That is the precondition for promoting any log-only "
            "policy to enforce.")
    else:
        # dimensionless source: shape only
        takeaways.append(
            f"Volume is the only lens available on {source} today. "
            f"Shape shifts in this timeline are the cleanest early "
            "warning for ingestion pipeline changes.")

    slide_chart(prs, "Volume and timing", eyebrow="section 2",
                chart_path=charts_dir / "02_daily_timeline.png",
                takeaways=takeaways[:3], accent=AMBER)


def slide_top_users(prs, d, charts_dir):
    summary = d["summary"]
    p_key = summary.get("top_principal_key")
    if not p_key or not summary.get("top_user"):
        return False
    tu = summary["top_user"]
    who = tu.get(p_key)
    label = _principal_label(p_key)
    top_n = int(tu.get("n") or 0)
    top_share = summary.get("top_share", 0.0)

    takeaways = [
        f"{who} = {top_n:,} events ({top_share:.1f}% of volume).",
    ]
    mix = _mix_rows(d)
    ranked = _derive_by_principal(mix, p_key)
    if len(ranked) >= 3:
        second = ranked[1][1]
        last = ranked[-1][1]
        takeaways.append(
            f"The rest of the cohort lives in the {last:,} to "
            f"{second:,} range. That is where normal {label} behaviour "
            "hides.")
    takeaways.append(
        "Log scale used deliberately so a dominant principal does not "
        "flatten the long tail.")

    slide_chart(prs, "Who is driving the traffic", eyebrow="section 3",
                chart_path=charts_dir / "03_top_users_log.png",
                takeaways=takeaways[:3], accent=NAVY)
    return True


def slide_blocks_bypass(prs, d, charts_dir):
    summary = d["summary"]
    dims = d.get("dims", {})
    p_key = summary.get("top_principal_key")
    if not dims.get("action") or not p_key:
        return
    label = _principal_label(p_key)
    mix = _mix_rows(d)

    # blocks
    block_rows = [r for r in mix
                  if str(r.get("action")) == "block"
                  and r.get(p_key) not in (None, "", "null")]
    block_rows.sort(key=lambda r: -int(r.get("n") or 0))

    takeaways: List[str] = []
    if block_rows:
        top = block_rows[0]
        who = top.get(p_key)
        n = int(top.get("n") or 0)
        total_blocks = summary.get("block", 0)
        share = 100.0 * n / total_blocks if total_blocks else 0
        takeaways.append(
            f"{who} blocks = {n:,}. "
            f"That is {share:.1f}% of all blocks in the window.")
    if len(block_rows) >= 2:
        cohort = block_rows[1:]
        lo = min(int(r["n"]) for r in cohort)
        hi = max(int(r["n"]) for r in cohort)
        takeaways.append(
            f"The rest of the blocked cohort sits between {lo:,} and "
            f"{hi:,} blocks per {label.rstrip('s')}. Policy tuning or "
            "1:1 coaching has the highest ROI here.")
    takeaways.append(
        "Anyone labelled 'anonymous' or 'null' in the chart warrants "
        "a tagging pass so blocks can be attributed.")

    slide_chart(prs, "Highest-risk prompt senders", eyebrow="section 4",
                chart_path=charts_dir / "04_blocks_by_user.png",
                takeaways=takeaways[:3], accent=CORAL)

    # bypass
    if summary.get("bypass", 0) <= 0:
        return
    bypass_rows = [r for r in mix
                   if str(r.get("action")) == "bypass"
                   and r.get(p_key) not in (None, "", "null")]
    bypass_rows.sort(key=lambda r: -int(r.get("n") or 0))
    bypass_total = summary.get("bypass", 0)
    bypass_takeaways = [
        f"{bypass_total:,} bypasses across the window. Low-volume, "
        "high-signal.",
    ]
    if bypass_rows:
        top2 = bypass_rows[:2]
        names = ", ".join(str(r.get(p_key)) for r in top2)
        bypass_takeaways.append(
            f"Top overriders: {names}. Cross-check endpoint posture "
            "for each.")
    bypass_takeaways.append(_bypass_short(summary.get("bypass_pct", 0.0)))

    slide_chart(prs, "Who is overriding guardrails", eyebrow="section 5",
                chart_path=charts_dir / "05_bypass_by_user.png",
                takeaways=bypass_takeaways[:3], accent=TEAL)


def slide_user_action_mix(prs, d, charts_dir):
    summary = d["summary"]
    dims = d.get("dims", {})
    p_key = summary.get("top_principal_key")
    if not p_key or not dims.get("action"):
        return
    label = _principal_label(p_key)

    mix = _mix_rows(d)
    # find the user with the most even split (mixed-use profile)
    by_prin: Dict[str, Dict[str, int]] = defaultdict(lambda: defaultdict(int))
    for r in mix:
        who = r.get(p_key)
        if who in (None, "", "null"):
            continue
        by_prin[str(who)][str(r.get("action"))] += int(r.get("n") or 0)
    mixed = []
    for who, acts in by_prin.items():
        tot = sum(acts.values())
        if tot < 5:
            continue
        share = max(acts.values()) / tot
        mixed.append((share, who, acts))
    mixed.sort(key=lambda x: x[0])

    takeaways = [
        f"Stacked action mix per top {label}. This is the fastest way "
        "to spot a mixed-use profile (blocked and allowed inside one "
        "session)."
    ]
    if mixed:
        share, who, acts = mixed[0]
        actions = ", ".join(f"{k}:{v}"
                            for k, v in sorted(acts.items(),
                                               key=lambda kv: -kv[1]))
        takeaways.append(
            f"{who} shows the clearest mixed-use profile "
            f"({actions}). Likely a real analyst probing the "
            "guardrails.")
    takeaways.append(
        "Use this chart to pick principals for policy feedback loops "
        "before the next re-run.")

    slide_chart(prs, "Per-principal action mix", eyebrow="section 6",
                chart_path=charts_dir / "06_user_action_mix.png",
                takeaways=takeaways[:3], accent=AMBER)


def slide_tenant_context(prs, d, charts_dir):
    summary = d["summary"]
    source = d["source"]
    rank = summary.get("rank_24h")
    takeaways = []
    if rank:
        takeaways.append(
            f"In the last 24 hours, {source} ranks #{rank} among the "
            "data sources landing in this tenant.")
    takeaways.append(
        f"Same data lake as SentinelOne EDR, identity, network, and "
        "cloud telemetry. Correlation is a PowerQuery join away.")
    takeaways.append(
        f"As usage grows, {source}'s share of daily events climbs in "
        "rank with zero ingestion changes.")

    slide_chart(prs,
                f"Where {source} sits in the broader tenant",
                eyebrow="section 7",
                chart_path=charts_dir / "07_tenant_context.png",
                takeaways=takeaways[:3], accent=SLATE)


def slide_recommendations(prs, d):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title_bar(s, "Recommendations, prioritised",
                    eyebrow="what to do with these findings")

    items = _recommendations(d)[:4]
    n = len(items)
    start_left = Inches(0.5)
    start_top = Inches(1.55)
    gap_x = Inches(0.3)
    gap_y = Inches(0.3)

    # layout adapts to count so the slide never half-empties
    if n >= 3:
        cols = 2
        rows = (n + 1) // 2
        card_w = Inches(6.0)
        card_h = Inches(2.35) if rows == 2 else Inches(4.95)
    elif n == 2:
        cols = 2
        rows = 1
        card_w = Inches(6.0)
        card_h = Inches(4.95)
    else:  # n == 1
        cols = 1
        rows = 1
        card_w = Inches(12.3)
        card_h = Inches(4.95)

    for idx, (when, color, title, body) in enumerate(items):
        row = idx // cols
        col = idx % cols
        left = start_left + col * (card_w + gap_x)
        top = start_top + row * (card_h + gap_y)
        card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  left, top, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(0xDD, 0xE1, 0xE6)
        card.line.width = Pt(0.5)
        card.shadow.inherit = False
        add_accent_bar(s, left, top, Inches(0.15), card_h, color)
        add_text(s, left + Inches(0.35), top + Inches(0.25),
                 card_w - Inches(0.5), Inches(0.3),
                 when, size=10, bold=True, color=color)
        add_text(s, left + Inches(0.35), top + Inches(0.55),
                 card_w - Inches(0.5), Inches(0.5),
                 title, size=18, bold=True, color=NAVY)
        # body height grows with card height
        body_h = card_h - Inches(1.1) - Inches(0.2)
        add_text(s, left + Inches(0.35), top + Inches(1.1),
                 card_w - Inches(0.5), body_h,
                 body, size=13 if card_h > Inches(3) else 12,
                 color=INK, line_spacing=1.3)


def slide_closing(prs, d):
    summary = d["summary"]
    source = d["source"]
    p_key = summary.get("top_principal_key")
    tu = summary.get("top_user")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, NAVY)
    add_accent_bar(s, 0, Inches(2.8), Inches(0.4), Inches(2.8), CORAL)

    add_text(s, Inches(0.9), Inches(2.7),
             Inches(11), Inches(0.4),
             "NEXT STEPS",
             size=14, bold=True, color=CORAL)

    # stitch a headline from the data
    if tu and summary.get("top_share", 0.0) >= 40 and p_key:
        headline = (f"Resolve {tu.get(p_key)}, promote\n"
                    "log-only policy, enrich ingestion.")
    elif (d.get("dims", {}).get("action")
          and summary.get("block_pct", 0) >= 10
          and summary.get("bypass_pct", 0) < 1):
        headline = ("Promote log-only policy, correlate\n"
                    "bypass, enrich ingestion.")
    else:
        headline = (f"Baseline {source}, correlate with\n"
                    "endpoint, and re-run this report.")

    add_text(s, Inches(0.9), Inches(3.1),
             Inches(11.5), Inches(1.4),
             headline,
             size=36, bold=True, color=WHITE, line_spacing=1.1)

    add_text(s, Inches(0.9), Inches(5.1),
             Inches(11), Inches(0.8),
             "Everything in this deck is reproducible from a single "
             "JSON artefact.\nRe-run any time the workload shifts.",
             size=15, italic=True, color=LIGHT, line_spacing=1.3)

    add_text(s, Inches(0.9), Inches(6.8),
             Inches(11), Inches(0.3),
             "Thank you. Questions?",
             size=14, bold=True, color=LIGHT)


# ------------------------------------------------------- main

def build_deck(data_path: Path) -> Path:
    d = json.loads(data_path.read_text())
    charts_dir = data_path.parent / "charts"
    slug = d.get("slug") or _slugify(d.get("source", "source"))
    window = d.get("window_label", "window")
    source = d["source"]
    out_path = data_path.parent / f"{slug}_CTO_Deck_{window}.pptx"

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_cover(prs, d)
    slide_execsummary(prs, d)
    slide_action_mix(prs, d, charts_dir)
    slide_daily_timeline(prs, d, charts_dir)
    slide_top_users(prs, d, charts_dir)
    slide_blocks_bypass(prs, d, charts_dir)
    slide_user_action_mix(prs, d, charts_dir)
    slide_tenant_context(prs, d, charts_dir)
    slide_recommendations(prs, d)
    slide_closing(prs, d)

    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, start=1):
        if i == 1 or i == total:
            continue
        add_footer(slide, i, total, source, window)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"wrote {out_path}  ({out_path.stat().st_size:,} bytes, "
          f"{total} slides)")
    return out_path


def main():
    ap = argparse.ArgumentParser(description=(
        "Build a CTO-facing PPTX deck from a data.json artefact."))
    ap.add_argument("--data", required=True,
                    help="Path to reports/<slug>_<window>/data.json")
    args = ap.parse_args()
    build_deck(Path(args.data))


if __name__ == "__main__":
    main()
